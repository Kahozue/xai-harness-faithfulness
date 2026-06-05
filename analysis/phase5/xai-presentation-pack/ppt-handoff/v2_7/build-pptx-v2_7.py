#!/usr/bin/env python3
"""把 v2_7 HTML 的每一頁渲染圖各自鋪滿一張 16:9 投影片，組成 .pptx。

來源：preview/slide-01.png .. slide-28.png（1920x1080，與 v2_7 HTML 同步渲染）。
輸出：ppt-handoff/xAI-faithfulness-harness-v2_7.pptx
忠實保留 HTML 版面，不重排文字。
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches

HERE = Path(__file__).resolve().parent          # .../ppt-handoff/v2_7
PREVIEW = HERE / "preview"
OUT = HERE.parent / "xAI-faithfulness-harness-v2_7.pptx"  # .../ppt-handoff/

# 16:9，1920x1080 對應 13.333in x 7.5in，整頁鋪滿不變形。
SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

slides = sorted(PREVIEW.glob("slide-*.png"))
if len(slides) != 28:
    raise SystemExit(f"expected 28 slide PNGs, found {len(slides)}")

prs = Presentation()
prs.slide_width = SLIDE_W
prs.slide_height = SLIDE_H
blank = prs.slide_layouts[6]  # 全空白版面

for png in slides:
    slide = prs.slides.add_slide(blank)
    slide.shapes.add_picture(str(png), 0, 0, width=SLIDE_W, height=SLIDE_H)

prs.save(str(OUT))
print(f"wrote {OUT} with {len(prs.slides._sldIdLst)} slides")
